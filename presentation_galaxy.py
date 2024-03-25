from io import BytesIO
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import requests
from picture import search_google_images

def create_slides_galaxy(gpt_response, slide_title, slide_number, Topic, choice, pic):
    prs = Presentation()
    #Topic Slide
    topic = prs.slides.add_slide(prs.slide_layouts[5])
    background = topic.shapes.add_picture("Templates\Galaxy\Galaxy_Title_Slide.png", 0, 0, prs.slide_width, prs.slide_height)
    #Topic Title
    t_title = topic.shapes.add_textbox(Pt(235), Pt(170), Pt(250), Pt(180))
    tf1 = t_title.text_frame
    t_title.text = Topic
    tf1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf1.paragraphs[0].font.size = Pt(54)
    tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf1.word_wrap = True
    #Topic Subtitle
    t_content = topic.shapes.add_textbox(Pt(240), Pt(415), Pt(240), Pt(30))
    ctf1 = t_content.text_frame
    t_content.text = "Created by SlideCraft"
    ctf1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ctf1.paragraphs[0].font.size = Pt(26)
    ctf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf1.word_wrap = True
    no_of_pic = 0
    if pic == 1:
        no_of_pic = slide_number/3
    count = 0
    for i in range(slide_number):
        slide_content = []
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        #Formatting Content from GPT
        mid_content = gpt_response[i].split(': ', 1)[-1].strip()
        if choice == 1:
            slide_content = str(mid_content)
        elif choice == 2:
            print(gpt_response)
            print("\n\n")
            split_data = gpt_response[i].split('\n')
            print(split_data)
            print("\n\n")
            if slide_number == 6 or slide_number == 9:
                mid_data = split_data[0].split('-')
                slide_content.append(mid_data[0].split(':')[1].strip())
                slide_content.append(mid_data[1].split(':')[1].strip())
            elif slide_number == 3:
                slide_content.append(split_data[0].split(':')[1].strip())
                slide_content.append(split_data[1].split(':')[1].strip())
        #slide content
        if count<no_of_pic:
            if i%2==0:
                background = slide.shapes.add_picture("Templates\Galaxy\Galaxy_Slide_"+str(i+1)+".png", 0, 0, prs.slide_width, prs.slide_height)
                content = slide.shapes.add_textbox(Pt(135), Pt(190), Pt(450), Pt(240))
                ctf2 = content.text_frame
                if choice == 1:
                    content.text = slide_content
                elif choice == 2:
                    content.text = str(slide_content[1])
                ctf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                ctf2.paragraphs[0].font.size = Pt(22)
                ctf2.paragraphs[0].alignment = PP_ALIGN.CENTER
                ctf2.word_wrap = True
                print("\n\nSlide done\n\n")
            else:
                count = count + 1
                background = slide.shapes.add_picture("Templates\Galaxy\Galaxy_Picture_Slide_"+str(count)+".png", 0, 0, prs.slide_width, prs.slide_height)
                content = slide.shapes.add_textbox(Pt(60), Pt(200), Pt(370), Pt(300))
                ctf2 = content.text_frame
                if choice == 1:
                    content.text = slide_content
                elif choice == 2:
                    content.text = str(slide_content[1])
                ctf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                ctf2.paragraphs[0].font.size = Pt(22)
                ctf2.paragraphs[0].alignment = PP_ALIGN.CENTER
                ctf2.word_wrap = True
                print("\n Slide "+str(i)+" Done \n")
                #Picture
                image_url = search_google_images(Topic+" and "+slide_content[0])
                image = str(image_url).strip("[]'")
                try:
                    response = requests.get(image)
                    image_stream = BytesIO(response.content)
                    picture = slide.shapes.add_picture(image_stream, Pt(480), Pt(245), Pt(190), Pt(135))
                    print("\n Photo "+count+" Done \n")
                except Exception as e:
                    print(f"Error inserting image {image_url}: {e}")
                print("\n\nSlide done\n\n")
        else:
            background = slide.shapes.add_picture("Templates\Galaxy\Galaxy_Slide_"+str(i+1)+".png", 0, 0, prs.slide_width, prs.slide_height)
            content = slide.shapes.add_textbox(Pt(135), Pt(190), Pt(450), Pt(240))
            ctf2 = content.text_frame
            if choice == 1:
                content.text = slide_content
            elif choice == 2:
                content.text = str(slide_content[1])
            ctf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            ctf2.paragraphs[0].font.size = Pt(22)
            ctf2.paragraphs[0].alignment = PP_ALIGN.CENTER
            ctf2.word_wrap = True
            print("\n\nSlide done\n\n")
        #slide title
        title = slide.shapes.add_textbox(Pt(110), Pt(40), Pt(500), Pt(70))
        tf2 = title.text_frame
        if choice == 1:
            title.text = slide_title[i].upper()
        elif choice == 2:
            title.text = str(slide_content[0])
        tf2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf2.paragraphs[0].font.size = Pt(50)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf2.word_wrap = True
        print("\n\nSlide done\n\n")

    prs.save(Topic+".pptx")
    print("\nPresentation Created Successfully.")
    os.system(f'start "" "{Topic+".pptx"}"')