from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def create_slides_landscape(gpt_response, slide_title, slide_number, Topic):
    prs = Presentation()
    #Topic Slide
    topic = prs.slides.add_slide(prs.slide_layouts[5])
    background = topic.shapes.add_picture("Templates\Landscape\Landscape_Topic_Slide.png", 0, 0, prs.slide_width, prs.slide_height)
    #Topic Title
    t_title = topic.shapes.add_textbox(Pt(235), Pt(170), Pt(250), Pt(180))
    tf1 = t_title.text_frame
    t_title.text = Topic
    tf1.paragraphs[0].font.color.rgb = RGBColor(237, 74, 5)
    tf1.paragraphs[0].font.size = Pt(48)
    tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf1.word_wrap = True
    #Topic Subtitle
    t_content = topic.shapes.add_textbox(Pt(240), Pt(340), Pt(240), Pt(30))
    ctf1 = t_content.text_frame
    t_content.text = "Created by ppt generator"
    ctf1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ctf1.paragraphs[0].font.size = Pt(22)
    ctf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf1.word_wrap = True
    #Content Slides
    for i in range(slide_number):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  
        background = slide.shapes.add_picture("Templates\Landscape\Landscape_Slide_"+str(i+1)+".png", 0, 0, prs.slide_width, prs.slide_height)
        #Content Slide Title
        title = slide.shapes.add_textbox(Pt(110), Pt(60), Pt(500), Pt(70))
        tf2 = title.text_frame
        title.text = slide_title[i].upper()
        tf2.paragraphs[0].font.color.rgb = RGBColor(3, 0, 97)
        tf2.paragraphs[0].font.size = Pt(44)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf2.word_wrap = True
        #Formatting Content from GPT
        gpt_response[i] = gpt_response[i].split(': ', 1)[-1].strip()
        slide_content = str(gpt_response[i])
        #Content Slide content
        content = slide.shapes.add_textbox(Pt(135), Pt(190), Pt(450), Pt(240))
        ctf2 = content.text_frame
        content.text = slide_content
        ctf2.paragraphs[0].font.color.rgb = RGBColor(255, 176, 48)
        ctf2.paragraphs[0].font.size = Pt(18)
        ctf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        ctf2.word_wrap = True

    prs.save(Topic+".pptx")
    print("\nPresentation Created Successfully.")